import json
import re
import requests
from pptx import Presentation
from pptx.dml.color import RGBColor

ANTHROPIC_API_KEY = "" # paste anthropic key
PPT_FILE = "input/Py-Slides-1.pptx"
MAX_KEYWORDS = 5

# Colors for highlighting - i think we can change the colors, based on preference, this is just a sample
RED = RGBColor(255, 0, 0)
GREEN = RGBColor(0, 150, 0)

def extract_slide_text(ppt_path):
    """extract text from each slide"""
    prs = Presentation(ppt_path)
    slides_text = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text + " "
        slides_text.append({"slide_num": slide_num, "text": text.strip()})
    
    return slides_text, prs

def get_keywords_from_api(slide_text, slide_num):
    """get keywords from API with specific content-reading prompt"""
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        "Content-Type": "application/json",
        "x-api-key": ANTHROPIC_API_KEY,
        "anthropic-version": "2023-06-01"
    }
    
    prompt = f"""Read this slide content carefully. Do NOT assume keywords based on domain name.

Extract exactly {MAX_KEYWORDS} UNIQUE, IMPORTANT study keywords that:
1. Actually appear in this specific slide text
2. Are key concepts someone would highlight while studying  
3. Are NOT common words (like 'python', 'language', 'program', 'code')
4. Are NOT duplicated in your response
5. Focus on technical terms, specific concepts, or important details
6. Must be words/phrases that actually exist in the slide content

Slide content: {slide_text}

Return ONLY JSON format: {{"keywords": ["word1", "word2", "word3", "word4", "word5"]}}"""
    
    data = {
        "model": "claude-3-sonnet-20240229",
        "max_tokens": 200,
        "messages": [{"role": "user", "content": prompt}]
    }
    
    try:
        response = requests.post(url, headers=headers, json=data)
        if response.status_code == 200:
            result = response.json()
            response_text = result["content"][0]["text"].strip()
            
            # Extract JSON from response
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                data = json.loads(json_match.group())
                keywords = data.get("keywords", [])
                
                # Validate keywords exist in slide text
                validated_keywords = []
                for keyword in keywords:
                    if keyword and re.search(r'\b' + re.escape(keyword) + r'\b', slide_text, re.IGNORECASE):
                        validated_keywords.append(keyword)
                
                return validated_keywords
            
        print(f"API error for slide {slide_num}: {response.status_code}")
        return []
        
    except Exception as e:
        print(f"Error getting keywords for slide {slide_num}: {e}")
        return []

def highlight_keywords(slide, keywords, slide_num):
    """highlight keywords in a slide with duplicate prevention"""
    if not keywords:
        return 0
    
    highlights = 0
    colors = [RED, GREEN]
    highlighted_words = set()  # track highlighted words to prevent duplicates
    
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                # get paragraph text
                para_text = ""
                for run in paragraph.runs:
                    para_text += run.text
                
                if not para_text.strip():
                    continue
                
                # find keywords in paragraph (avoid duplicates)
                found_keywords = []
                for i, keyword in enumerate(keywords):
                    if keyword.lower() not in highlighted_words:
                        pattern = r'\b' + re.escape(keyword) + r'\b'
                        if re.search(pattern, para_text, re.IGNORECASE):
                            found_keywords.append((keyword, i))
                            highlighted_words.add(keyword.lower())
                
                if found_keywords:
                    # rebuild paragraph with highlights
                    paragraph.clear()
                    remaining_text = para_text
                    
                    # get all keyword positions
                    positions = []
                    for keyword, color_idx in found_keywords:
                        pattern = r'\b' + re.escape(keyword) + r'\b'
                        for match in re.finditer(pattern, remaining_text, re.IGNORECASE):
                            positions.append({
                                'start': match.start(),
                                'end': match.end(),
                                'text': match.group(),
                                'color': colors[color_idx % 2]
                            })
                    
                    # remove overlapping positions
                    positions.sort(key=lambda x: x['start'])
                    filtered_positions = []
                    for pos in positions:
                        overlap = False
                        for existing in filtered_positions:
                            if (pos['start'] < existing['end'] and pos['end'] > existing['start']):
                                overlap = True
                                break
                        if not overlap:
                            filtered_positions.append(pos)
                    
                    # rebuild paragraph
                    last_pos = 0
                    for pos in filtered_positions:
                        # add text before keyword
                        if pos['start'] > last_pos:
                            run = paragraph.add_run()
                            run.text = remaining_text[last_pos:pos['start']]
                        
                        # add highlighted keyword
                        run = paragraph.add_run()
                        run.text = pos['text']
                        run.font.color.rgb = pos['color']
                        run.font.bold = True
                        highlights += 1
                        
                        last_pos = pos['end']
                    
                    # add remaining text
                    if last_pos < len(remaining_text):
                        run = paragraph.add_run()
                        run.text = remaining_text[last_pos:]
    
    return highlights

def main():
    """main function"""
    print(f"Processing: {PPT_FILE}")
    
    # extract slide text
    slides_text, prs = extract_slide_text(PPT_FILE)
    print(f"Found {len(slides_text)} slides")
    
    total_highlights = 0
    
    # process each slide
    for slide_data in slides_text:
        slide_num = slide_data["slide_num"]
        text = slide_data["text"]
        
        if not text:
            print(f"Slide {slide_num}: No text")
            continue
        
        # get keywords from API
        print(f"Slide {slide_num}: Getting keywords...")
        keywords = get_keywords_from_api(text, slide_num)
        
        print(f"Slide {slide_num}: Keywords={keywords}")
        
        # highlight keywords
        if keywords:
            slide = prs.slides[slide_num - 1]
            highlights = highlight_keywords(slide, keywords, slide_num)
            total_highlights += highlights
            print(f"Slide {slide_num}: Applied {highlights} highlights")
    
    # save result
    output_path = PPT_FILE.replace('.pptx', '_highlighted.pptx')
    prs.save(output_path)
    
    print(f"Complete! Total highlights: {total_highlights}")
    print(f"Saved to: {output_path}")

if __name__ == "__main__":
    main()